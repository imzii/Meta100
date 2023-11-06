import re
import zlib
import struct
import olefile
import sys
from tkinter import Tk
from tkinter.filedialog import askopenfilename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

def get_hwp_text(filename):
    f = olefile.OleFileIO(filename)
    dirs = f.listdir()

    # HWP 파일 검증
    if ['FileHeader'] not in dirs or \
            ['\x05HwpSummaryInformation'] not in dirs:
        raise Exception('Not Valid HWP.')

    # 문서 포맷 압축 여부 확인
    header = f.openstream('FileHeader')
    header_data = header.read()
    is_compressed = (header_data[36] & 1) == 1

    # Body Sections 불러오기
    nums = []
    for d in dirs:
        if d[0] == 'BodyText':
            nums.append(int(d[1][len('Section'):]))
    sections = ['BodyText/Section' + str(x) for x in sorted(nums)]

    # 전체 text 추출
    text = ''
    for section in sections:
        body_text = f.openstream(section)
        data = body_text.read()
        if is_compressed:
            unpacked_data = zlib.decompress(data, -15)
        else:
            unpacked_data = data

        # 각 Section 내 text 추출
        section_text = ''
        i = 0
        size = len(unpacked_data)
        while i < size:
            header = struct.unpack_from('<I', unpacked_data, i)[0]
            rec_type = header & 0x3ff
            rec_len = (header >> 20) & 0xfff

            if rec_type in [67]:
                rec_data = unpacked_data[i + 4:i + 4 + rec_len]
                section_text += rec_data.decode('utf-16')
                section_text += '\n'

            i += 4 + rec_len

        text += section_text
        text += '\n'

    return text


def remove_stopwords(text):
    # 정규 표현식 패턴을 사용하여 텍스트에서 불필요한 문자 제거
    pattern = re.compile(r'[^A-Za-z0-9가-힣\s.?!,~\'"”’#$%&*()-:;\x0B\x0D]')
    text = re.sub(pattern, '', text)
    text = text.replace('\u001F\u001F', '-')
    return text

def draw_background(slide, slide_width, slide_height):
    # 슬라이드 배경색 설정
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검정색 배경

    # 배경 이미지 추가
    background_path = 'src/background.jpg'
    background_left = 0
    background_top = 0
    background_width = slide_width
    background_height = slide_height
    slide.shapes.add_picture(background_path, background_left, background_top, background_width, background_height)

def draw_logo(slide):
    # 로고 이미지 추가
    logo_path = 'src/logo.png'
    logo_left = Inches(9.1)
    logo_top = Inches(0.05)
    logo_width = Inches(0.78)
    logo_height = Inches(0.48)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, logo_width, logo_height)

def create_slide(prs, with_background, with_logo):
    # 슬라이드 생성 함수
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃
    slide = prs.slides.add_slide(slide_layout)
    if with_background:
        draw_background(slide, slide_width, slide_height)
    if with_logo:
        draw_logo(slide)
    return slide

def create_title_slide(prs, title):
    # 타이틀 슬라이드 생성 함수
    slide = create_slide(prs, True, False)

    # 텍스트 상자 설정
    text_box_margin = Inches(1)
    title_textbox = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=text_box_margin / 2,
        top=slide_height / 4,
        width=slide_width - text_box_margin,
        height=slide_height / 2,
    )
    title_textbox.fill.background()
    title_textbox.line.fill.background()

    text_frame = title_textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.name = 'HY강M'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글자

def create_slide_with_text_and_shape(prs, text, unit):
    # 슬라이드 생성 함수 (텍스트 및 도형 포함)
    slide = create_slide(prs, True, True)

    text_box_margin = Inches(0.2)

    # 투명 직사각형 도형 추가
    rectangle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=text_box_margin / 2,
        top=slide_height / 4,
        width=slide_width - text_box_margin,
        height=slide_height / 2,
    )
    rectangle.fill.background()
    rectangle.line.fill.background()

    text_frame = rectangle.text_frame
    text_frame.word_wrap = True
    text_frame.line_spacing = Pt(48)

    p = text_frame.paragraphs[0]
    p.line_spacing = Pt(48)
    p.text = text
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    p.font.name = '맑은 고딕'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글자

    # 도형 추가
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,  # 직사각형 도형
        left=Inches(0.23),
        top=Inches(0.15),
        width=Inches(1.55),
        height=Inches(0.44)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(148, 0, 211)  # 보라색 배경색

    shape.line.color.rgb = RGBColor(255, 165, 0)  # 주황색 윤곽색
    shape.line.width = Pt(4.5)

    text_frame = shape.text_frame
    text_frame.word_wrap = True

    # 텍스트 스타일 설정
    p = text_frame.paragraphs[0]
    p.text = unit
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.name = 'HY견고딕'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글자

def split_text_by_passages(text):
    pattern = r'\n+\d+\.\s*|\n+\d+\s*~\s*\d+\.\s*'
    matches = re.finditer(pattern, text)
    indices = [match.start() for match in matches]

    sentences = []
    start_index = 0
    for index in indices:
        sentences.append(text[start_index:index].strip())
        start_index = index

    sentences.append(text[start_index:].strip())

    return sentences


def split_text_by_sentences(text):
    pattern = r'(?<=[.!?])\s+'
    sentences = re.split(pattern, text)
    return [sentence.strip() for sentence in sentences if sentence.strip()]

if __name__ == '__main__':
    if len(sys.argv) != 2:
        print('Usage: HwpToPpt.py [title]')
        sys.exit()

    title = sys.argv[1]

    Tk().withdraw()
    hwp_file_name = askopenfilename()
    hwp_file_text = remove_stopwords(get_hwp_text(hwp_file_name))

    pptx = Presentation()
    slide_height = pptx.slide_height
    slide_width = pptx.slide_width

    create_title_slide(pptx, title)

    passages = split_text_by_passages(hwp_file_text)

    current_unit = ''
    for passage in passages:
        sentences = split_text_by_sentences(passage)
        for index in range(len(sentences)):
            pattern = r'\b(\d+\.\s*|\d+\s*~\s*\d+\.\s*)'
            if index < len(sentences) - 1 and sentences[index + 1][0] == '*':
                sentences[index + 1] = sentences[index] + '\n\n' + sentences[index + 1]
                continue
            if re.fullmatch(pattern, sentences[index]):
                current_unit = sentences[index].replace('.', '번')
            else:
                create_slide_with_text_and_shape(pptx, sentences[index], current_unit)

    output_pptx_filename = hwp_file_name.replace('Desktop', 'Desktop/Presentation').replace('hwp', 'pptx')
    pptx.save(output_pptx_filename)
    print('저장 완료')